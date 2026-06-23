"""
Build a Nature-style Chinese journal club PPTX from slide config JSON.
Usage: python build_pptx.py <slides.json> <figures_dir> <output.pptx>
"""
import os, sys, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def build_pptx(slides_json_path, figures_dir, output_path):
    with open(slides_json_path, 'r', encoding='utf-8') as f:
        slides_data = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Nature-style palette
    DB = RGBColor(0x1B, 0x3A, 0x5C)
    AB = RGBColor(0x3A, 0x7C, 0xA5)
    LG = RGBColor(0xF2, 0xF2, 0xF2)
    DG = RGBColor(0x33, 0x33, 0x33)
    MG = RGBColor(0x66, 0x66, 0x66)
    WH = RGBColor(0xFF, 0xFF, 0xFF)

    def bg(s, c=WH):
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = c

    def TB(s, l, t, w, h, tx, fs=18, c=DG, b=False, a=PP_ALIGN.LEFT):
        bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        bx.text_frame.word_wrap = True
        p = bx.text_frame.paragraphs[0]
        p.text = tx; p.font.size = Pt(fs); p.font.color.rgb = c
        p.font.bold = b; p.font.name = 'Arial'; p.alignment = a

    def BL(s, l, t, w, h, items, fs=16, c=DG, ls=1.5):
        bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        bx.text_frame.word_wrap = True
        for i, item in enumerate(items):
            p = bx.text_frame.paragraphs[0] if i == 0 else bx.text_frame.add_paragraph()
            p.text = item; p.font.size = Pt(fs); p.font.color.rgb = c
            p.font.name = 'Arial'; p.space_after = Pt(fs * (ls - 1) * 0.7)

    def LN(s, l, t, w, c=AB, lw=Pt(2)):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), lw)
        sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

    def RC(s, l, t, w, h, c):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

    def FT(s, tx=''):
        LN(s, 0.8, 7.1, 11.733, AB, Pt(0.8))
        TB(s, 0.8, 7.15, 5, 0.3, tx, fs=9, c=MG)

    def SR(s, tx, l=0.8, t=7.05):
        TB(s, l, t, 8, 0.25, tx, fs=8, c=MG)

    def PIC(s, fn, l, t, w, h):
        p = os.path.join(figures_dir, fn)
        if os.path.exists(p):
            s.shapes.add_picture(p, Inches(l), Inches(t), Inches(w), Inches(h))

    for sd in slides_data:
        st = sd['type']

        if st == 'cover':
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bg(s); RC(s, 0, 0, 13.333, 0.08, AB); RC(s, 0, 7.42, 13.333, 0.08, AB)
            TB(s, 1.0, 1.5, 11.333, 2.0, sd['title'], fs=36, c=DB, b=True)
            LN(s, 1.0, 3.8, 3.0, AB, Pt(3))
            TB(s, 1.0, 4.1, 11.333, 0.8, sd.get('subtitle', ''), fs=18, c=MG)
            TB(s, 1.0, 5.3, 11.333, 0.5, sd.get('authors', ''), fs=14, c=MG)
            TB(s, 1.0, 5.8, 11.333, 0.4, sd.get('journal', ''), fs=12, c=MG)
            TB(s, 1.0, 6.3, 11.333, 0.4, sd.get('affil', ''), fs=12, c=MG)

        elif st == 'bullets':
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bg(s); RC(s, 0, 0, 13.333, 0.06, AB)
            TB(s, 0.8, 0.4, 11.733, 0.6, sd['title'], fs=30, c=DB, b=True)
            LN(s, 0.8, 1.05, 2.5, AB, Pt(2.5))
            fs_val = sd.get('font_size', 16)
            BL(s, 0.8, 1.4, 11.733, 5.5, sd['bullets'], fs=fs_val, c=DG, ls=1.3)
            FT(s, sd.get('footer', ''))

        elif st == 'figure':
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bg(s); RC(s, 0, 0, 13.333, 0.06, AB)
            TB(s, 0.8, 0.4, 11.733, 0.6, sd['title'], fs=26, c=DB, b=True)
            LN(s, 0.8, 1.0, 2.5, AB, Pt(2.5))
            PIC(s, sd['figure'], 0.8, 1.3, 11.733, 5.5)
            SR(s, sd.get('source', ''))
            FT(s, sd.get('footer', ''))

        elif st == 'dual_panel':
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bg(s); RC(s, 0, 0, 13.333, 0.06, AB)
            TB(s, 0.8, 0.4, 11.733, 0.6, sd['title'], fs=28, c=DB, b=True)
            LN(s, 0.8, 1.05, 2.5, AB, Pt(2.5))
            L, R = sd['left'], sd['right']
            RC(s, 0.8, 1.5, 5.5, 2.8, LG)
            TB(s, 1.0, 1.6, 5.1, 0.5, L['header'], fs=20, c=DB, b=True)
            BL(s, 1.0, 2.2, 5.1, 2.0, L['bullets'], fs=14, c=DG, ls=1.2)
            RC(s, 6.8, 1.5, 5.5, 2.8, LG)
            TB(s, 7.0, 1.6, 5.1, 0.5, R['header'], fs=20, c=DB, b=True)
            BL(s, 7.0, 2.2, 5.1, 2.0, R['bullets'], fs=14, c=DG, ls=1.2)
            if sd.get('summary'):
                BL(s, 0.8, 4.8, 11.733, 2.0, sd['summary'], fs=17, c=DB, ls=1.6)
            FT(s, sd.get('footer', ''))

        elif st == 'three_columns':
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bg(s); RC(s, 0, 0, 13.333, 0.06, AB)
            TB(s, 0.8, 0.4, 11.733, 0.6, sd['title'], fs=26, c=DB, b=True)
            LN(s, 0.8, 1.0, 2.5, AB, Pt(2.5))
            for i, col in enumerate(sd['cols']):
                x = 0.8 + i * 4.0
                RC(s, x, 1.4, 3.6, 5.3, LG)
                TB(s, x + 0.1, 1.5, 3.4, 0.4, col['header'], fs=20, c=DB, b=True)
                BL(s, x + 0.1, 2.0, 3.4, 4.5, col['bullets'], fs=13, c=DG, ls=1.3)
            FT(s, sd.get('footer', ''))

        elif st == 'dual_with_figure':
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bg(s); RC(s, 0, 0, 13.333, 0.06, AB)
            TB(s, 0.8, 0.4, 11.733, 0.6, sd['title'], fs=28, c=DB, b=True)
            LN(s, 0.8, 1.0, 2.5, AB, Pt(2.5))
            L, R = sd['left'], sd['right']
            RC(s, 0.8, 1.3, 5.5, 2.8, LG)
            TB(s, 1.0, 1.4, 5.1, 0.5, L['header'], fs=20, c=DB, b=True)
            BL(s, 1.0, 2.0, 5.1, 2.0, L['bullets'], fs=13, c=DG, ls=1.25)
            RC(s, 6.8, 1.3, 5.5, 2.8, LG)
            TB(s, 7.0, 1.4, 5.1, 0.5, R['header'], fs=20, c=DB, b=True)
            BL(s, 7.0, 2.0, 5.1, 2.0, R['bullets'], fs=13, c=DG, ls=1.25)
            PIC(s, sd['figure'], 1.5, 4.3, 10.333, 2.6)
            SR(s, sd.get('source', ''))
            FT(s, sd.get('footer', ''))

        elif st == 'innovations':
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bg(s); RC(s, 0, 0, 13.333, 0.06, AB)
            TB(s, 0.8, 0.4, 11.733, 0.6, sd['title'], fs=30, c=DB, b=True)
            LN(s, 0.8, 1.05, 2.5, AB, Pt(2.5))
            for i, (hdr, desc) in enumerate(sd['items']):
                y = 1.4 + i * 1.4
                RC(s, 0.8, y, 11.733, 1.2, LG)
                TB(s, 1.0, y + 0.1, 5.0, 0.4, hdr, fs=20, c=DB, b=True)
                TB(s, 1.0, y + 0.5, 11.3, 0.6, desc, fs=15, c=DG)
            FT(s, sd.get('footer', ''))

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return len(prs.slides)


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python build_pptx.py <slides.json> <figures_dir> <output.pptx>")
        sys.exit(1)
    n = build_pptx(sys.argv[1], sys.argv[2], sys.argv[3])
    print(f"PPTX saved: {sys.argv[3]} ({n} slides)")
